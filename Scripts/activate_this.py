import tkinter as tk
from tkinter import filedialog
from tkinter import ttk
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from docx2pdf import convert as convert_word_to_pdf
from PIL import Image


def open_file_dialog(event):
    # Reset the conversion status label
    conversion_status_label.config(text="")
    selected_input_type = input_type_var.get()
    # Define filetypes based on the selected input type
    filetypes = [("PPT Files", "*.ppt;*.pptx")] if selected_input_type == "PPT File" \
        else [("Word Files", "*.doc;*.docx")] if selected_input_type == "Word File" \
        else [("Image Files", "*.png;*.jpg;*.jpeg")]
    file_paths = filedialog.askopenfilenames(title=f"Select {selected_input_type} files", filetypes=filetypes)
    print(f"Selected {selected_input_type} files:", file_paths)
    # Set the selected file paths in the entry widget (optional)
    selected_file_entry.delete("1.0", tk.END)
    for file_path in file_paths:
        selected_file_entry.insert(tk.END, file_path + "\n")


def start_conversion():
    input_file_type = input_type_var.get()
    # Retrieve all the text in the Text widget
    input_file_paths = selected_file_entry.get("1.0", tk.END).split('\n')
    # Remove the last empty string from the list
    if input_file_paths[-1] == '':
        input_file_paths.pop()
    if not input_file_paths:
        print("Please select files.")
        return
    # Use the save as dialog to get the output PDF file path
    output_pdf_file = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF Files", "*.pdf")])
    if not output_pdf_file:
        print("Conversion canceled.")
        return
    try:
        # Perform the conversion based on the selected input type
        if input_file_type == "Word File":
            convert_word_to_pdf(input_file_paths[0], output_pdf_file)
        elif input_file_type == "PPT File":
            ppt_to_pdf(input_file_paths[0], output_pdf_file)
        elif input_file_type == "Image File":
            # Pass a list of image file paths
            convert_images_to_pdf(input_file_paths, output_pdf_file)
        conversion_status_label.config(text="Conversion Completed Successfully", fg="green")
        print("Conversion successful.")
    except Exception as e:
        conversion_status_label.config(text=f"Conversion failed. Error: {e}", fg="red")
        print(f"Conversion failed. Error: {e}")


def extract_text_from_shape(shape):
    if shape.has_text_frame:
        return shape.text_frame.text
    elif shape.has_table:
        table_text = ""
        for row in shape.table.rows:
            for cell in row.cells:
                table_text += cell.text + " "
            table_text += "\n"
        return table_text
    elif shape.has_chart:
        return "Chart: " + shape.chart_title.text_frame.text
    else:
        return ""

def ppt_to_pdf(input_ppt, output_pdf):
    presentation = Presentation(input_ppt)
    # Create a PDF document
    pdf_doc = SimpleDocTemplate(output_pdf, pagesize=letter)
    flowables = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        flowables.append(Paragraph(f"<b>Slide {slide_number}</b>", getSampleStyleSheet()['Heading1']))
        for shape in slide.shapes:
            text = extract_text_from_shape(shape)
            if text:
                # Add a spacer to simulate line break
                flowables.append(Spacer(1, 12))
                # Add the extracted text with a custom style
                flowables.append(Paragraph(text, ParagraphStyle('Normal')))
    # Build the PDF document
    pdf_doc.build(flowables)


def convert_images_to_pdf(image_paths, output_pdf):
    c = canvas.Canvas(output_pdf, pagesize=letter)
    for image_path in image_paths:
        img = Image.open(image_path)
        pdf_width, pdf_height = letter
        pdf_aspect = pdf_width / pdf_height
        img_width, img_height = img.size
        img_aspect = img_width / img_height
        if img_aspect > pdf_aspect:
            img_width = pdf_width
            img_height = int(pdf_width / img_aspect)
        else:
            img_height = pdf_height
            img_width = int(pdf_height * img_aspect)
        x_offset = (pdf_width - img_width) / 2
        y_offset = (pdf_height - img_height) / 2
        c.drawImage(image_path, x_offset, y_offset, width=img_width, height=img_height)
        c.showPage()
    c.save()


# Create the main window
root = tk.Tk()
root.title("Effortless Easy PDF Convertor")
# Set the size of the window (width x height)
root.geometry("500x300")
# Add a heading with blue text color
heading_label = tk.Label(root, text="Effortless Easy PDF Convertor", font=("Arial", 24, "bold"), fg="blue")
heading_label.pack(pady=10)
# Create a Combobox (dropdown list) with increased size
input_types = ["PPT File", "Word File", "Image File"]
input_type_var = tk.StringVar()
input_type_dropdown = ttk.Combobox(root, textvariable=input_type_var, values=input_types, font=("Arial", 14))
input_type_dropdown.set("Select Input Type")
input_type_dropdown.pack(pady=10)
# Bind the event to the open_file_dialog function
input_type_dropdown.bind("<<ComboboxSelected>>", open_file_dialog)
# Text widget to display the selected file paths
selected_file_entry = tk.Text(root, height=5, width=50)
selected_file_entry.pack(pady=5)
# Add a button for starting the conversion with increased size
start_button = tk.Button(root, text="Start The Conversion", command=start_conversion, font=("Arial", 16))
start_button.pack(pady=10)
# Label to display conversion status
conversion_status_label = tk.Label(root, text="", font=("Arial", 12))
conversion_status_label.pack(pady=5)
# Start the main loop
root.mainloop()
